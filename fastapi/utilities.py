import json
import re
import os
from typing import List, Dict, Any
try:
    from openai import OpenAI
    import PyPDF2
except ImportError:
    import PyPDF2
    print("ERROR: Missing necessary dependencies, please run:")
    print("pip install openai PyPDF2")
    exit(1)

def parse_llm_json(response: str) -> dict:
    # Strip Markdown code fences if present
    response = response.strip()
    response = re.sub(r"^```(?:json)?\s*", "", response)
    response = re.sub(r"\s*```$", "", response)
    response = response.strip()

    if not response:
        raise ValueError("LLM returned an empty response")

    return json.loads(response)

def extract_text_from_pdf(pdf_path: str, document_title: str = None) -> str:
    try:
        with open(pdf_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            title = document_title or os.path.basename(pdf_path)
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n--- Page {page_num} | Source: {title} ---\n\n"
                    text += page_text
            return text
    except Exception as e:
        raise Exception(f"Failed to extract text from {pdf_path}: {e}")

def extract_text_from_html(html_path: str, document_title: str = None) -> str:
    try:
        from bs4 import BeautifulSoup

        with open(html_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")

        title = document_title or os.path.basename(html_path)

        # Remove non-content elements
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()

        text = ""
        current_section = []
        section_num = 1

        for element in soup.find_all(["h1", "h2", "h3", "h4", "p", "pre", "code", "li", "table"]):
            if element.name in ("h1", "h2", "h3"):
                if current_section:
                    section_text = " ".join(current_section).strip()
                    if section_text:
                        text += f"\n\n--- Page {section_num} | Source: {title} ---\n\n"
                        text += section_text
                        section_num += 1
                current_section = [element.get_text(strip=True)]
            else:
                content = element.get_text(strip=True)
                if content:
                    current_section.append(content)

        # Last section
        if current_section:
            section_text = " ".join(current_section).strip()
            if section_text:
                text += f"\n\n--- Page {section_num} | Source: {title} ---\n\n"
                text += section_text

        return text
    except ImportError:
        raise Exception("beautifulsoup4 not installed. Run: pip install beautifulsoup4")
    except Exception as e:
        raise Exception(f"Failed to extract text from {html_path}: {e}")


def extract_text_from_pptx(pptx_path: str, document_title: str = None) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        title = document_title or os.path.basename(pptx_path)
        text = ""

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            slide_text.append(para_text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            slide_text.append(row_text)

            # Speaker notes often explain what graphs show
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[Notes: {notes}]")

            if slide_text:
                text += f"\n\n--- Page {slide_num} | Source: {title} ---\n\n"
                text += "\n".join(slide_text)

        return text
    except ImportError:
        raise Exception("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        raise Exception(f"Failed to extract text from {pptx_path}: {e}")


def extract_text_from_source(path: str, document_title: str = None) -> str:
    """Unified text extraction, dispatches based on file extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path, document_title)
    elif ext in (".html", ".htm"):
        return extract_text_from_html(path, document_title)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(path, document_title)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .html, .htm, .pptx")

# Splitting extracted text into sections based on pages, for now
def split_into_sections(text: str) -> List[Dict[str, Any]]:
    """Split text into sections, preserving page number and source title."""
    raw_sections = text.split("--- Page")
    sections = []
    for s in raw_sections:
        s = s.strip()
        if not s:
            continue
        # Parse header line: "42 | Source: Practical C Programming ---"
        lines = s.split("\n", 1)
        header = lines[0]
        content = lines[1].strip() if len(lines) > 1 else ""

        page_num = None
        source = None
        try:
            # "42 | Source: Practical C Programming ---"
            parts = header.split("|")
            page_num = int(parts[0].strip())
            source = parts[1].replace("Source:", "").replace("---", "").strip()
        except Exception:
            pass

        sections.append({
            "page": page_num,
            "source": source,
            "content": content
        })
    return sections


def parse_json(text: str) -> Any:
    """Attempts to parse JSON string into Python object"""
    try:
        # Strip markdown fences if model ignored instructions
        text = re.sub(r"^```(?:json)?\s*", "", text.strip())
        text = re.sub(r"\s*```$", "", text).strip()

        return json.loads(text)
    except json.JSONDecodeError:
        # Second attempt: try to fix invalid escape sequences
        try:
            # Replace invalid backslash escapes that aren't valid JSON escapes
            # Valid JSON escapes are: \" \\ \/ \b \f \n \r \t \uXXXX
            cleaned = re.sub(r'\\(?!["\\/bfnrtu])', r'\\\\', text)
            return json.loads(cleaned)
        except json.JSONDecodeError:
            print("\n=== RAW MODEL OUTPUT START ===")
            print(text)
            print("=== RAW MODEL OUTPUT END ===\n")
            raise

def load_json(path: str) -> List[Dict[str, Any]]:
        try:
            with open(path, "r", encoding="utf-8") as f:
                print(f"[Import] Successfully loaded {path}")
                return json.load(f)
        except Exception as e:
            raise Exception(f"[Import] Failed to load from {path}: {e}")


def export_json(units: List[Dict[str, Any]], output_path: str) -> None:
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(units, f, indent=2, ensure_ascii=False)
        print(f"[Export] Successfully saved {len(units)} units to {output_path}")
    except Exception as e:
        print(f"[Export] Failed to save to {output_path}: {e}")